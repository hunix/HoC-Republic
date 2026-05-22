#!/usr/bin/env python3
"""
HoC Professional Presentation Generator v2
Templates: executive, modern-light, gradient, minimal, corporate
Layouts: title, section, content, two_column, image_text, chart, comparison,
         closing, stats, timeline, quote, table, team, process, swot
"""
import json, os, sys, datetime, textwrap, io, urllib.request, subprocess, re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Parse args from JSON file ───────────────────────────────────
_args_path = "/tmp/_gen_args.json"
if os.path.exists(_args_path):
    with open(_args_path, "r") as f:
        args = json.load(f)
elif len(sys.argv) > 1:
    args = json.loads(sys.argv[1])
else:
    args = {}
slide_data_raw = args.get("slide_data", "[]")
branding_raw = args.get("branding", "{}")
title_text = args.get("title", "Untitled")
out_path = args.get("out_path", "/workspace/presentation.pptx")
template_name = args.get("template", "executive")

try:    parsed = json.loads(slide_data_raw)
except: parsed = []

if isinstance(parsed, dict):
    slides_list = parsed.get("slides", [])
    brand_data = parsed.get("branding", {})
else:
    slides_list = parsed if isinstance(parsed, list) else []
    brand_data = {}

try:
    brand_override = json.loads(branding_raw)
    if brand_override: brand_data = {**brand_data, **brand_override}
except: pass

# ── SVG→PNG Conversion ─────────────────────────────────────────
def svg_to_png(svg_path, png_path=None, width=800):
    """Convert SVG to PNG using CairoSVG or Inkscape fallback."""
    if not png_path:
        png_path = svg_path.rsplit(".", 1)[0] + ".png"
    try:
        import cairosvg
        cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=width)
        return png_path
    except ImportError:
        pass
    try:
        from PIL import Image
        import subprocess as sp
        r = sp.run(["inkscape", svg_path, "--export-type=png",
                     f"--export-filename={png_path}", f"-w{width}"],
                    capture_output=True, timeout=15)
        if r.returncode == 0 and os.path.exists(png_path):
            return png_path
    except: pass
    try:
        from PIL import Image
        r = subprocess.run(["rsvg-convert", "-w", str(width), svg_path, "-o", png_path],
                           capture_output=True, timeout=10)
        if r.returncode == 0: return png_path
    except: pass
    return None

def download_image(url, dest):
    """Download an image, auto-convert SVG to PNG."""
    try:
        urllib.request.urlretrieve(url, dest)
        # Detect SVG by extension or content
        is_svg = dest.lower().endswith(".svg")
        if not is_svg:
            with open(dest, "rb") as f:
                header = f.read(256)
                is_svg = b"<svg" in header or b"<!DOCTYPE svg" in header
        if is_svg:
            png_dest = dest.rsplit(".", 1)[0] + ".png"
            converted = svg_to_png(dest, png_dest)
            if converted: return converted
            # Pillow fallback for simple SVGs
            try:
                subprocess.run(["pip", "install", "-q", "cairosvg"], capture_output=True, timeout=30)
                import cairosvg
                cairosvg.svg2png(url=dest, write_to=png_dest, output_width=800)
                return png_dest
            except: pass
        return dest
    except:
        return None

# ── Template Color Palettes ─────────────────────────────────────
def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

TEMPLATES = {
    "executive": {
        "bg_primary":   "#0f172a", "bg_secondary": "#1e293b",
        "accent":       "#3b82f6", "accent2":      "#8b5cf6",
        "text_light":   "#ffffff", "text_muted":   "#94a3b8",
        "text_dark":    "#0f172a", "success":      "#22c55e",
        "warning":      "#f59e0b", "danger":       "#ef4444",
        "gradient_start": "#1e3a5f", "gradient_end": "#0f172a",
    },
    "modern-light": {
        "bg_primary":   "#ffffff", "bg_secondary": "#f8fafc",
        "accent":       "#2563eb", "accent2":      "#7c3aed",
        "text_light":   "#ffffff", "text_muted":   "#64748b",
        "text_dark":    "#0f172a", "success":      "#16a34a",
        "warning":      "#d97706", "danger":       "#dc2626",
        "gradient_start": "#eff6ff", "gradient_end": "#ffffff",
    },
    "gradient": {
        "bg_primary":   "#1a1a2e", "bg_secondary": "#16213e",
        "accent":       "#e94560", "accent2":      "#f97316",
        "text_light":   "#ffffff", "text_muted":   "#a0aec0",
        "text_dark":    "#1a1a2e", "success":      "#4ade80",
        "warning":      "#fbbf24", "danger":       "#f87171",
        "gradient_start": "#e94560", "gradient_end": "#f97316",
    },
    "minimal": {
        "bg_primary":   "#fafafa", "bg_secondary": "#f5f5f5",
        "accent":       "#18181b", "accent2":      "#71717a",
        "text_light":   "#ffffff", "text_muted":   "#a1a1aa",
        "text_dark":    "#18181b", "success":      "#22c55e",
        "warning":      "#eab308", "danger":       "#ef4444",
        "gradient_start": "#f4f4f5", "gradient_end": "#fafafa",
    },
    "corporate": {
        "bg_primary":   "#1e3a5f", "bg_secondary": "#2d4a6f",
        "accent":       "#0ea5e9", "accent2":      "#06b6d4",
        "text_light":   "#ffffff", "text_muted":   "#93c5fd",
        "text_dark":    "#0c1929", "success":      "#34d399",
        "warning":      "#fbbf24", "danger":       "#f87171",
        "gradient_start": "#1e3a5f", "gradient_end": "#0c2340",
    },
}

# Apply branding overrides on top of template
tpl = TEMPLATES.get(template_name, TEMPLATES["executive"])
PRIMARY   = hex_to_rgb(brand_data.get("primary_color", tpl["bg_primary"]))
SECONDARY = hex_to_rgb(brand_data.get("secondary_color", tpl["bg_secondary"]))
ACCENT    = hex_to_rgb(brand_data.get("accent_color", tpl["accent"]))
ACCENT2   = hex_to_rgb(tpl["accent2"])
TEXT_LIGHT = hex_to_rgb(tpl["text_light"])
TEXT_MUTED = hex_to_rgb(tpl["text_muted"])
TEXT_DARK  = hex_to_rgb(tpl["text_dark"])
SUCCESS   = hex_to_rgb(tpl["success"])
WARNING   = hex_to_rgb(tpl["warning"])
DANGER    = hex_to_rgb(tpl["danger"])
GRAD_START = hex_to_rgb(tpl["gradient_start"])
GRAD_END   = hex_to_rgb(tpl["gradient_end"])
IS_DARK = template_name in ("executive", "gradient", "corporate")
BODY_TEXT = TEXT_LIGHT if IS_DARK else TEXT_DARK

COMPANY = brand_data.get("company_name", "")
LOGO_URL = brand_data.get("logo_url", "")
FONT = brand_data.get("font_family", "Calibri")
TODAY = datetime.date.today().strftime("%B %d, %Y")

# ── Presentation Setup ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = 13.333
SH = 7.5

# ── Helper Functions ────────────────────────────────────────────
def set_bg(slide, color=PRIMARY):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_text(slide, text, left, top, w, h, size=18, color=BODY_TEXT, bold=False, align=PP_ALIGN.LEFT, font_name=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]; p.text = str(text); p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    p.font.name = font_name or FONT
    return box

def add_rich_text(slide, text, left, top, w, h, size=16, color=BODY_TEXT, line_spacing=1.2):
    """Add text with bullet-point support and markdown-lite formatting."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    lines = str(text).split("\\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        stripped = line.lstrip("- •*▸▹►")
        is_bullet = line.lstrip().startswith(("-", "•", "*", "▸", "►"))
        is_sub = line.startswith("  ") and is_bullet
        if is_sub:
            p.text = "    ◦  " + stripped.strip()
            p.font.size = Pt(size - 2)
        elif is_bullet:
            p.text = "  •  " + stripped.strip()
            p.font.size = Pt(size)
        else:
            # Check for **bold** inline
            p.text = line
            p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(int(size * 0.3))
        p.line_spacing = Pt(int(size * line_spacing))
    return box

def add_header_bar(slide, title_str, subtitle=""):
    """Accent-colored header bar at top of content slides."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(slide, title_str, 0.6, 0.12, SW-1.2, 0.5, size=28, color=TEXT_LIGHT, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.6, 0.55, SW-1.2, 0.3, size=12, color=hex_to_rgb("#cbd5e1"))

def add_footer(slide, idx, total):
    """Footer with slide number, company, and date."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(SH-0.5), Inches(SW-1), Inches(0.015))
    line.fill.solid(); line.fill.fore_color.rgb = TEXT_MUTED; line.line.fill.background()
    fy = SH - 0.4
    if COMPANY:
        add_text(slide, COMPANY, 0.6, fy, 4, 0.3, size=9, color=TEXT_MUTED)
    add_text(slide, TODAY, SW/2-1.5, fy, 3, 0.3, size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, f"{idx+1} / {total}", SW-2, fy, 1.4, 0.3, size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

def add_card(slide, left, top, w, h, fill_color=SECONDARY, shadow=False):
    """Rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if shadow:
        from pptx.oxml.ns import qn
        # Simple shadow via shape offset
        pass
    return shape

def add_icon_text(slide, icon, text, left, top, size=14):
    """Icon + text combo for inline elements."""
    add_text(slide, f"{icon}  {text}", left, top, SW-2, 0.4, size=size, color=BODY_TEXT)

def generate_chart(chart_type, chart_data, chart_path):
    """Generate a chart image using matplotlib."""
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.ticker as mticker

        fig, ax = plt.subplots(figsize=(8, 4.5))
        bg = tpl["bg_primary"] if IS_DARK else "#ffffff"
        fg = tpl["bg_secondary"] if IS_DARK else "#f0f0f0"
        fig.patch.set_facecolor(bg); ax.set_facecolor(fg)

        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        colors_list = chart_data.get("colors", None)
        default_colors = [tpl["accent"], tpl["accent2"], tpl["success"], tpl["warning"], tpl["danger"],
                          "#06b6d4", "#f97316", "#ec4899"]

        tc = "white" if IS_DARK else "#333333"

        if chart_type == "pie":
            c = colors_list or default_colors[:len(labels)]
            wedges, texts, autotexts = ax.pie(values, labels=labels, colors=c, autopct="%1.1f%%",
                textprops={"color": tc, "fontsize": 11}, startangle=90,
                wedgeprops={"edgecolor": bg, "linewidth": 2})
            for t in texts: t.set_color(tc)
        elif chart_type == "line":
            ax.plot(labels, values, color=tpl["accent"], linewidth=2.5, marker="o", markersize=8)
            ax.fill_between(range(len(values)), values, alpha=0.15, color=tpl["accent"])
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels)
        elif chart_type == "horizontal_bar":
            c = colors_list or [tpl["accent"]]*len(labels)
            bars = ax.barh(labels, values, color=c, height=0.5, edgecolor="none")
            for bar, v in zip(bars, values):
                ax.text(bar.get_width() + max(values)*0.01, bar.get_y() + bar.get_height()/2,
                        f"{v:,.0f}" if isinstance(v, (int,float)) else str(v),
                        va="center", color=tc, fontsize=10)
        elif chart_type == "donut":
            c = colors_list or default_colors[:len(labels)]
            wedges, texts, autotexts = ax.pie(values, labels=labels, colors=c, autopct="%1.1f%%",
                textprops={"color": tc, "fontsize": 10}, startangle=90, pctdistance=0.8,
                wedgeprops={"width": 0.4, "edgecolor": bg, "linewidth": 2})
        else:  # bar
            c = colors_list or [tpl["accent"]]*len(labels)
            bars = ax.bar(labels, values, color=c, width=0.6, edgecolor="none")
            for bar, v in zip(bars, values):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(values)*0.01,
                        f"{v:,.0f}" if isinstance(v, (int,float)) else str(v),
                        ha="center", va="bottom", color=tc, fontsize=10, fontweight="bold")

        if chart_type not in ("pie", "donut"):
            ax.tick_params(colors=tc, labelsize=10)
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
            sc = "#475569" if IS_DARK else "#d1d5db"
            ax.spines["bottom"].set_color(sc); ax.spines["left"].set_color(sc)
            ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format(int(x), ",")))

        ct = chart_data.get("title", "")
        if ct: ax.set_title(ct, color=tc, fontsize=14, fontweight="bold", pad=12)

        plt.tight_layout()
        fig.savefig(chart_path, dpi=200, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig)
        return True
    except Exception as e:
        print(f"Chart generation failed: {e}", file=sys.stderr)
        return False

def try_download_logo():
    if not LOGO_URL: return None
    return download_image(LOGO_URL, "/tmp/_brand_logo.png")

logo_path = try_download_logo()

# ── Build Slides ────────────────────────────────────────────────
total = len(slides_list) or 1
if not slides_list:
    slides_list = [{"layout": "title", "title": title_text, "subtitle": "Generated by HoC"}]
    total = 1

for i, s in enumerate(slides_list):
    layout = s.get("layout", "content")
    st = s.get("title", f"Slide {i+1}")
    sc = s.get("content", "")
    ssub = s.get("subtitle", "")
    snotes = s.get("notes", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── TITLE SLIDE ─────────────────────────────────────────────
    if layout == "title" or (i == 0 and layout == "content" and not sc):
        set_bg(slide, PRIMARY)
        # Decorative accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), Inches(SH))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()
        # Decorative accent circle (modern touch)
        if template_name in ("gradient", "modern-light"):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW-3), Inches(-1), Inches(4), Inches(4))
            circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT2; circle.line.fill.background()
            circle.rotation = 0
        # Logo
        if logo_path and os.path.exists(logo_path):
            try: slide.shapes.add_picture(logo_path, Inches(1), Inches(0.8), height=Inches(0.8))
            except: pass
        add_text(slide, st, 1, 2.2, 10, 1.8, size=48, color=TEXT_LIGHT if IS_DARK else TEXT_DARK, bold=True)
        if ssub or sc:
            add_text(slide, ssub or sc, 1, 4.2, 9, 1, size=22, color=TEXT_MUTED)
        if COMPANY:
            add_text(slide, COMPANY, 1, 5.8, 6, 0.5, size=14, color=ACCENT)
        add_text(slide, TODAY, 1, 6.4, 4, 0.4, size=12, color=TEXT_MUTED)

    # ── SECTION DIVIDER ─────────────────────────────────────────
    elif layout == "section":
        set_bg(slide, ACCENT)
        # Decorative element
        deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SW/2-2), Inches(2.0), Inches(4), Inches(0.06))
        deco.fill.solid(); deco.fill.fore_color.rgb = TEXT_LIGHT; deco.line.fill.background()
        add_text(slide, st, 1, 2.5, SW-2, 2, size=44, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)
        if sc:
            add_text(slide, sc, 2, 4.8, SW-4, 1, size=20, color=hex_to_rgb("#e2e8f0"), align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── CONTENT SLIDE ───────────────────────────────────────────
    elif layout == "content":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st, ssub)
        add_rich_text(slide, sc, 0.6, 1.2, SW-1.2, SH-2, size=16, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── TWO COLUMN ──────────────────────────────────────────────
    elif layout == "two_column":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        parts = sc.split("|||") if "|||" in str(sc) else [str(sc), ""]
        col_w = (SW - 1.6) / 2 - 0.2
        add_card(slide, 0.6, 1.2, col_w, SH-2.0)
        add_rich_text(slide, parts[0], 0.9, 1.4, col_w-0.6, SH-2.4, size=14, color=BODY_TEXT)
        add_card(slide, 0.6+col_w+0.4, 1.2, col_w, SH-2.0)
        add_rich_text(slide, parts[1] if len(parts) > 1 else "", 0.6+col_w+0.7, 1.4, col_w-0.6, SH-2.4, size=14, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── IMAGE + TEXT ────────────────────────────────────────────
    elif layout == "image_text":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        img_url = s.get("image_url", "")
        img_path = download_image(img_url, f"/tmp/_slide_img_{i}.png") if img_url else None
        half_w = (SW - 1.2) / 2
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.2), width=Inches(half_w-0.2), height=Inches(SH-2.0))
            except:
                add_text(slide, "[Image could not be loaded]", 0.6, 3, half_w, 1, size=14, color=TEXT_MUTED)
            add_rich_text(slide, sc, 0.6+half_w+0.2, 1.2, half_w-0.2, SH-2.0, size=15, color=BODY_TEXT)
        else:
            add_rich_text(slide, sc, 0.6, 1.2, SW-1.2, SH-2.0, size=16, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── CHART SLIDE ─────────────────────────────────────────────
    elif layout == "chart":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        ct = s.get("chart_type", "bar")
        cd = s.get("chart_data", {})
        cp = f"/tmp/_chart_{i}.png"
        if generate_chart(ct, cd, cp):
            try: slide.shapes.add_picture(cp, Inches(1.5), Inches(1.3), width=Inches(10), height=Inches(5))
            except: add_text(slide, "[Chart render failed]", 2, 3, 8, 1, size=16, color=WARNING)
        else:
            add_text(slide, "[Chart data insufficient]", 2, 3, 8, 1, size=16, color=WARNING)
        if snotes:
            add_text(slide, snotes, 0.6, 6.3, SW-1.2, 0.5, size=11, color=TEXT_MUTED)
        add_footer(slide, i, total)

    # ── COMPARISON TABLE ────────────────────────────────────────
    elif layout == "comparison":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        columns = s.get("columns", [])
        if columns:
            num_cols = len(columns)
            col_w = min((SW - 1.2) / num_cols - 0.3, 4.0)
            start_x = (SW - (col_w + 0.3) * num_cols + 0.3) / 2
            for ci, col in enumerate(columns):
                cx = start_x + ci * (col_w + 0.3)
                add_card(slide, cx, 1.2, col_w, 0.7, ACCENT)
                add_text(slide, col.get("header", ""), cx+0.2, 1.3, col_w-0.4, 0.5, size=16, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)
                for ri, item in enumerate(col.get("items", [])):
                    ry = 2.1 + ri * 0.55
                    bg_c = SECONDARY if ri % 2 == 0 else PRIMARY
                    add_card(slide, cx, ry, col_w, 0.5, bg_c)
                    add_text(slide, str(item), cx+0.2, ry+0.05, col_w-0.4, 0.4, size=13, color=BODY_TEXT, align=PP_ALIGN.CENTER)
        else:
            add_rich_text(slide, sc, 0.6, 1.2, SW-1.2, SH-2.0, size=16, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── STATS / KPI CARDS ───────────────────────────────────────
    elif layout == "stats":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        stats = s.get("stats", [])
        if not stats and sc:
            # Parse from content: "Revenue: $1.2M | Users: 50K | Growth: +23%"
            for part in sc.split("|"):
                kv = part.strip().split(":")
                if len(kv) == 2:
                    stats.append({"label": kv[0].strip(), "value": kv[1].strip()})
        num = len(stats) or 1
        card_w = min((SW - 1.2) / num - 0.3, 3.5)
        sx = (SW - (card_w + 0.3) * num + 0.3) / 2
        for si, stat in enumerate(stats):
            cx = sx + si * (card_w + 0.3)
            add_card(slide, cx, 1.6, card_w, 3.8, SECONDARY)
            # KPI value (large)
            val = stat.get("value", "—")
            add_text(slide, val, cx+0.2, 2.2, card_w-0.4, 1.2, size=42, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            # Label
            add_text(slide, stat.get("label", ""), cx+0.2, 3.6, card_w-0.4, 0.5, size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
            # Trend indicator
            trend = stat.get("trend", "")
            if trend:
                tc = SUCCESS if trend.startswith("+") or trend.startswith("↑") else DANGER if trend.startswith("-") or trend.startswith("↓") else TEXT_MUTED
                add_text(slide, trend, cx+0.2, 4.3, card_w-0.4, 0.5, size=14, color=tc, align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── TIMELINE ────────────────────────────────────────────────
    elif layout == "timeline":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        milestones = s.get("milestones", [])
        if not milestones and sc:
            for part in sc.split("|"):
                kv = part.strip().split(":")
                if len(kv) >= 2:
                    milestones.append({"date": kv[0].strip(), "title": ":".join(kv[1:]).strip()})
        nm = len(milestones) or 1
        # Horizontal timeline line
        ly = 3.5
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(ly), Inches(SW-1.6), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
        spacing = (SW - 2) / nm
        for mi, ms in enumerate(milestones):
            mx = 1.0 + mi * spacing
            # Dot on timeline
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(mx + spacing/2 - 0.15), Inches(ly - 0.12), Inches(0.28), Inches(0.28))
            dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
            # Date above
            add_text(slide, ms.get("date", ""), mx, ly - 1.0, spacing, 0.6, size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            # Title below
            add_text(slide, ms.get("title", ""), mx, ly + 0.4, spacing, 1.5, size=13, color=BODY_TEXT, align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── QUOTE SLIDE ─────────────────────────────────────────────
    elif layout == "quote":
        set_bg(slide, PRIMARY)
        # Large opening quote mark
        add_text(slide, "❝", 1, 1.0, 2, 2, size=120, color=ACCENT, bold=True)
        # Quote text
        add_text(slide, sc, 2, 2.5, SW-4, 3, size=28, color=BODY_TEXT, align=PP_ALIGN.CENTER)
        # Attribution
        attr = s.get("attribution", ssub or "")
        if attr:
            add_text(slide, f"— {attr}", 2, 5.5, SW-4, 0.8, size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── TABLE SLIDE ─────────────────────────────────────────────
    elif layout == "table":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        rows = s.get("rows", [])
        headers = s.get("headers", [])
        if not rows and sc:
            # Parse from pipe-delimited content
            for line in sc.split("\\n"):
                cells = [c.strip() for c in line.split("|") if c.strip()]
                if cells:
                    if not headers: headers = cells
                    else: rows.append(cells)
        if headers:
            ncols = len(headers)
            cw = min((SW - 1.2) / ncols, 3.0)
            sx = (SW - cw * ncols) / 2
            # Header row
            for ci, h in enumerate(headers):
                add_card(slide, sx + ci * cw, 1.2, cw - 0.05, 0.55, ACCENT)
                add_text(slide, h, sx + ci * cw + 0.1, 1.25, cw - 0.25, 0.45, size=14, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            # Data rows
            for ri, row in enumerate(rows[:8]):  # max 8 rows
                for ci, cell in enumerate(row[:ncols]):
                    ry = 1.8 + ri * 0.5
                    bg_c = SECONDARY if ri % 2 == 0 else PRIMARY
                    add_card(slide, sx + ci * cw, ry, cw - 0.05, 0.45, bg_c)
                    add_text(slide, str(cell), sx + ci * cw + 0.1, ry + 0.05, cw - 0.25, 0.35, size=12, color=BODY_TEXT, align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── TEAM SLIDE ──────────────────────────────────────────────
    elif layout == "team":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        members = s.get("members", [])
        if not members and sc:
            for part in sc.split("|"):
                kv = part.strip().split(":")
                if len(kv) >= 2:
                    members.append({"name": kv[0].strip(), "role": kv[1].strip()})
        nm = len(members) or 1
        cols = min(nm, 4)
        rows_count = (nm + cols - 1) // cols
        card_w = min((SW - 1.2) / cols - 0.3, 3.0)
        for mi, m in enumerate(members[:12]):
            r = mi // cols
            c = mi % cols
            sx = (SW - (card_w + 0.3) * cols + 0.3) / 2
            cx = sx + c * (card_w + 0.3)
            cy = 1.4 + r * 2.2
            add_card(slide, cx, cy, card_w, 1.9, SECONDARY)
            # Avatar circle placeholder
            avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + card_w/2 - 0.4), Inches(cy + 0.15), Inches(0.8), Inches(0.8))
            avatar.fill.solid(); avatar.fill.fore_color.rgb = ACCENT; avatar.line.fill.background()
            add_text(slide, m.get("name", ""), cx+0.1, cy+1.0, card_w-0.2, 0.4, size=14, color=BODY_TEXT, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, m.get("role", ""), cx+0.1, cy+1.35, card_w-0.2, 0.4, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        add_footer(slide, i, total)

    # ── PROCESS / FLOW SLIDE ────────────────────────────────────
    elif layout == "process":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        steps = s.get("steps", [])
        if not steps and sc:
            for idx_s, part in enumerate(sc.split("|")):
                steps.append({"title": part.strip(), "description": ""})
        ns = len(steps) or 1
        step_w = min((SW - 1.5) / ns - 0.2, 2.5)
        sx = (SW - (step_w + 0.4) * ns + 0.4) / 2
        for si, step in enumerate(steps[:6]):
            cx = sx + si * (step_w + 0.4)
            # Step number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + step_w/2 - 0.3), Inches(1.5), Inches(0.6), Inches(0.6))
            circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
            add_text(slide, str(si+1), cx + step_w/2 - 0.15, 1.55, 0.3, 0.5, size=18, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)
            # Arrow between steps
            if si < ns - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx + step_w + 0.05), Inches(1.75), Inches(0.3), Inches(0.04))
                arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_MUTED; arrow.line.fill.background()
            # Step card
            add_card(slide, cx, 2.4, step_w, 3.5, SECONDARY)
            add_text(slide, step.get("title", ""), cx+0.15, 2.6, step_w-0.3, 0.6, size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
            desc = step.get("description", "")
            if desc:
                add_rich_text(slide, desc, cx+0.15, 3.3, step_w-0.3, 2.4, size=11, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── SWOT / 2×2 MATRIX ──────────────────────────────────────
    elif layout == "swot":
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        quadrants = s.get("quadrants", [])
        if not quadrants and sc:
            parts = sc.split("|||")
            labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
            for qi, part in enumerate(parts[:4]):
                quadrants.append({"label": labels[qi] if qi < len(labels) else f"Q{qi+1}", "content": part.strip()})
        q_colors = [SUCCESS, DANGER, ACCENT, WARNING]
        positions = [(0.6, 1.2), (SW/2+0.15, 1.2), (0.6, SH/2+0.1), (SW/2+0.15, SH/2+0.1)]
        qw = SW/2 - 0.75
        qh = SH/2 - 1.0
        for qi, q in enumerate(quadrants[:4]):
            px, py = positions[qi]
            add_card(slide, px, py, qw, qh, SECONDARY)
            # Colored label bar at top of quadrant
            lbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px), Inches(py), Inches(qw), Inches(0.45))
            lbar.fill.solid(); lbar.fill.fore_color.rgb = q_colors[qi]; lbar.line.fill.background()
            add_text(slide, q.get("label", ""), px+0.2, py+0.05, qw-0.4, 0.35, size=14, color=TEXT_LIGHT, bold=True)
            add_rich_text(slide, q.get("content", ""), px+0.15, py+0.55, qw-0.3, qh-0.7, size=12, color=BODY_TEXT)
        add_footer(slide, i, total)

    # ── CLOSING SLIDE ───────────────────────────────────────────
    elif layout == "closing":
        set_bg(slide, PRIMARY)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(SH-0.15), Inches(SW), Inches(0.15))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()
        if logo_path and os.path.exists(logo_path):
            try: slide.shapes.add_picture(logo_path, Inches(SW/2-0.5), Inches(1.5), height=Inches(1.0))
            except: pass
        add_text(slide, st or "Thank You", 1, 3.0, SW-2, 1.5, size=44, color=TEXT_LIGHT if IS_DARK else TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        if sc:
            add_text(slide, sc, 2, 4.8, SW-4, 1, size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        if COMPANY:
            add_text(slide, COMPANY, 2, 6.0, SW-4, 0.5, size=14, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── FALLBACK ────────────────────────────────────────────────
    else:
        set_bg(slide, PRIMARY)
        add_header_bar(slide, st)
        add_rich_text(slide, sc, 0.6, 1.2, SW-1.2, SH-2.0, size=16, color=BODY_TEXT)
        add_footer(slide, i, total)

# ── Save ────────────────────────────────────────────────────────
prs.save(out_path)
print(f"Presentation saved: {out_path} ({total} slides, template: {template_name})")
